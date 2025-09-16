import streamlit as st
from openai import AzureOpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import requests
import io
import json
import os
from dotenv import load_dotenv
 
# Load environment variables from .env file
load_dotenv()

# ==============================
# CONFIG (using environment variables for security)
# ==============================
AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT", "")
AZURE_DEPLOYMENT_NAME = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME", "")
AZURE_API_KEY = os.getenv("AZURE_OPENAI_API_KEY", "")

# Check if required environment variables are set
if not all([AZURE_OPENAI_ENDPOINT, AZURE_DEPLOYMENT_NAME, AZURE_API_KEY]):
    st.error("⚠️ Missing Azure OpenAI configuration. Please check your .env file.")
    st.info("Required variables: AZURE_OPENAI_ENDPOINT, AZURE_OPENAI_DEPLOYMENT_NAME, AZURE_OPENAI_API_KEY")
    st.stop()
 
try:
    client = AzureOpenAI(
        api_key=AZURE_API_KEY,
        api_version="2024-05-01-preview",
        azure_endpoint=AZURE_OPENAI_ENDPOINT
    )
except Exception as e:
    st.error(f"❌ Failed to initialize Azure OpenAI client: {e}")
    st.stop()
 
# ==============================
# STREAMLIT UI
# ==============================
st.title("AI-Powered Presentation Generator (Green Theme)")
st.write("Upload a document and generate a green-branded presentation with AI (Azure OpenAI).")
st.info("ℹ️ Note: Image generation is currently disabled as your Azure OpenAI deployment doesn't include DALL-E.")
 
uploaded_file = st.file_uploader("Upload document (TXT/DOCX/PDF)", type=["txt", "docx", "pdf"])
 
if uploaded_file:
    content = ""
    if uploaded_file.type == "text/plain":
        content = uploaded_file.read().decode("utf-8")
    elif uploaded_file.type == "application/pdf":
        from PyPDF2 import PdfReader
        reader = PdfReader(uploaded_file)
        content = "\n".join([page.extract_text() for page in reader.pages])
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        import docx
        doc = docx.Document(uploaded_file)
        content = "\n".join([para.text for para in doc.paragraphs])
 
    st.subheader("Extracted Document Content")
    st.text_area("Document Content", content, height=200)
 
    if st.button("Generate Slides"):
        try:
            # ==============================
            # STEP 1: Ask LLM for structured outline
            # ==============================
            prompt = f"""
            You are an expert presentation creator. Based on the following content,
            create a slide deck outline. 
            
            IMPORTANT: Return ONLY valid JSON (no markdown formatting, no ```json blocks).
            
            Use this exact JSON structure:
            [
              {{
                "title": "Slide Title",
                "bullets": ["point 1", "point 2"],
                "image_prompt": "short text describing image to generate"
              }}
            ]
     
            Content:
            {content}
            
            Remember: Return pure JSON only, no markdown formatting.
            """
            response = client.chat.completions.create(
                model=AZURE_DEPLOYMENT_NAME,
                messages=[{"role": "user", "content": prompt}]
            )

            slide_plan = response.choices[0].message.content
            if not slide_plan:
                st.error("❌ Empty response from Azure OpenAI")
                st.stop()
            
            st.json(slide_plan)

            import json
            import re
            try:
                # Ensure slide_plan is not None before parsing
                if slide_plan is None:
                    st.error("❌ No response content from Azure OpenAI")
                    st.stop()
                
                # Type assertion to help linter understand slide_plan is not None
                assert slide_plan is not None
                
                # Clean up the response - remove markdown code blocks if present
                cleaned_response = slide_plan.strip()
                
                # Remove ```json and ``` markers if present
                if cleaned_response.startswith("```json"):
                    cleaned_response = cleaned_response[7:]  # Remove ```json
                elif cleaned_response.startswith("```"):
                    cleaned_response = cleaned_response[3:]   # Remove ```
                    
                if cleaned_response.endswith("```"):
                    cleaned_response = cleaned_response[:-3]  # Remove trailing ```
                    
                # Remove any remaining backticks or whitespace
                cleaned_response = cleaned_response.strip().strip('`')
                
                # Parse the cleaned JSON
                slides = json.loads(cleaned_response)
                if not isinstance(slides, list):
                    st.error("❌ Response is not a valid slide array")
                    st.stop()
            except json.JSONDecodeError as e:
                st.error(f"❌ Invalid JSON response: {e}")
                st.write("Raw response:", slide_plan)
                st.write("Cleaned response:", cleaned_response if 'cleaned_response' in locals() else "N/A")
                st.stop()
                
            # ==============================
            # STEP 2: Create PowerPoint (Green Theme)
            # ==============================
            st.info(f"📝 Creating presentation with {len(slides)} slides...")
            prs = Presentation()
            
            for i, slide in enumerate(slides):
                st.write(f"Creating slide {i+1}: {slide.get('title', 'Untitled')}")
                layout = prs.slide_layouts[5]  # Title + Content
                s = prs.slides.add_slide(layout)
     
                # Title formatting 🟢 (with safe null checking)
                title = s.shapes.title
                if title and hasattr(title, 'text'):
                    title.text = slide.get("title", "Untitled Slide")
                    if hasattr(title, 'text_frame') and title.text_frame:
                        title.text_frame.paragraphs[0].font.size = Pt(32)
                        title.text_frame.paragraphs[0].font.bold = True
                        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)  # Green
     
                # Bullet formatting 🟢 (with safe placeholder access)
                if len(s.placeholders) > 1:
                    content_placeholder = s.placeholders[1]
                    # Safe access to text_frame with getattr to avoid type issues
                    text_frame = getattr(content_placeholder, 'text_frame', None)
                    if text_frame:
                        try:
                            text_frame.clear()
                            bullets = slide.get("bullets", [])
                            for bullet in bullets:
                                p = text_frame.add_paragraph()
                                p.text = str(bullet)
                                p.level = 0
                                p.font.size = Pt(20)
                                p.font.color.rgb = RGBColor(0, 100, 0)  # Dark Green
                        except Exception as content_error:
                            st.warning(f"⚠️ Could not format content for slide '{slide.get('title', 'Unknown')}': {content_error}")
     
                # ==============================
                # STEP 3: Add image from DALL·E (with safe error handling)
                # Note: DALL-E requires separate deployment - skipping if not available
                # ==============================
                dalle_prompt = slide.get("image_prompt", "")
                if dalle_prompt and False:  # Temporarily disabled DALL-E due to deployment issues
                    try:
                        dalle_resp = client.images.generate(
                            model="dall-e-3",
                            prompt=dalle_prompt,
                            size="512x512"
                        )
                        if dalle_resp and dalle_resp.data and len(dalle_resp.data) > 0:
                            img_url = dalle_resp.data[0].url
                            if img_url:
                                img_response = requests.get(img_url)
                                if img_response.status_code == 200:
                                    image_stream = io.BytesIO(img_response.content)
                                    left = Inches(5.5)
                                    top = Inches(1.5)
                                    s.shapes.add_picture(image_stream, left, top, Inches(3), Inches(3))
                    except Exception as img_error:
                        st.warning(f"⚠️ Could not generate image for slide '{slide.get('title', 'Unknown')}': {img_error}")
                # Note: Image generation disabled - your Azure OpenAI deployment doesn't include DALL-E
     
                # Background branding 🟢
                try:
                    fill = s.background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(230, 255, 230)  # Light Green background
                except Exception as bg_error:
                    st.warning(f"⚠️ Could not set background color: {bg_error}")
     
            # Save presentation using in-memory buffer for better performance
            output_buffer = io.BytesIO()
            prs.save(output_buffer)
            output_buffer.seek(0)
            
            # Show presentation statistics
            presentation_size = len(output_buffer.getvalue())
            st.info(f"📊 Presentation created with {len(prs.slides)} slides ({presentation_size:,} bytes)")
     
            st.success("✅ Presentation generated successfully!")
            st.download_button(
                "Download Green-Branded Presentation", 
                output_buffer.getvalue(), 
                file_name="AI_Presentation_Green.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            
        except Exception as e:
            st.error(f"❌ Error generating presentation: {e}")
            st.write("Please check your Azure OpenAI configuration and try again.")